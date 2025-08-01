from fastapi import APIRouter, HTTPException
from fastapi.responses import FileResponse
import os
import tempfile
from datetime import datetime
from docx import Document
from docx.shared import Inches
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.units import inch
from pptx import Presentation
from pptx.util import Inches
from odf import text, teletype
from odf.opendocument import OpenDocumentText

router = APIRouter()

def create_temp_dir():
    """Cria um diretório temporário para os documentos"""
    temp_dir = tempfile.mkdtemp(prefix="fecitel_docs_")
    return temp_dir

def create_anais_doc():
    """Cria o documento Anais FECITEL.doc"""
    doc = Document()
    doc.add_heading('ANAIS FECITEL', 0)
    doc.add_paragraph('Feira de Ciência e Tecnologia')
    doc.add_paragraph(f'Data: {datetime.now().strftime("%d/%m/%Y")}')
    
    doc.add_heading('Sumário', level=1)
    doc.add_paragraph('1. Apresentação')
    doc.add_paragraph('2. Trabalhos Apresentados')
    doc.add_paragraph('3. Avaliações')
    doc.add_paragraph('4. Premiação')
    
    doc.add_heading('1. Apresentação', level=1)
    doc.add_paragraph('Este documento apresenta os anais da Feira de Ciência e Tecnologia (FECITEL), contendo todos os trabalhos apresentados, suas avaliações e resultados finais.')
    
    return doc

def create_apresentacao_pdf():
    """Cria o PDF da apresentação da feira"""
    temp_dir = create_temp_dir()
    filename = os.path.join(temp_dir, "apresentacao_feira.pdf")
    
    c = canvas.Canvas(filename, pagesize=A4)
    width, height = A4
    
    # Título
    c.setFont("Helvetica-Bold", 24)
    c.drawString(width/2 - 100, height - 100, "FECITEL")
    c.setFont("Helvetica", 16)
    c.drawString(width/2 - 80, height - 130, "Feira de Ciência e Tecnologia")
    
    # Informações
    c.setFont("Helvetica", 12)
    y_position = height - 200
    c.drawString(100, y_position, "• Evento anual de ciência e tecnologia")
    y_position -= 20
    c.drawString(100, y_position, "• Participação de escolas da região")
    y_position -= 20
    c.drawString(100, y_position, "• Avaliação por comissão especializada")
    y_position -= 20
    c.drawString(100, y_position, f"• Data: {datetime.now().strftime('%d/%m/%Y')}")
    
    c.save()
    return filename

def create_instrucoes_doc():
    """Cria o documento de instruções para avaliação"""
    doc = Document()
    doc.add_heading('INSTRUÇÕES PARA AVALIAÇÃO DE TRABALHOS NA FECITEL', 0)
    
    doc.add_heading('Critérios de Avaliação', level=1)
    doc.add_paragraph('1. Criatividade e Inovação (25 pontos)')
    doc.add_paragraph('2. Metodologia Científica (25 pontos)')
    doc.add_paragraph('3. Apresentação e Comunicação (20 pontos)')
    doc.add_paragraph('4. Relevância Social (15 pontos)')
    doc.add_paragraph('5. Viabilidade Técnica (15 pontos)')
    
    doc.add_heading('Procedimentos de Avaliação', level=1)
    doc.add_paragraph('• Cada avaliador deve preencher a ficha de avaliação')
    doc.add_paragraph('• A nota final será a média das avaliações')
    doc.add_paragraph('• Em caso de empate, será considerada a criatividade')
    
    return doc

def create_mensagem_avaliador_doc():
    """Cria a mensagem para o avaliador"""
    doc = Document()
    doc.add_heading('MENSAGEM AO AVALIADOR', 0)
    
    doc.add_paragraph('Prezado(a) Avaliador(a),')
    doc.add_paragraph('')
    doc.add_paragraph('Agradecemos sua participação na FECITEL - Feira de Ciência e Tecnologia.')
    doc.add_paragraph('')
    doc.add_paragraph('Sua avaliação é fundamental para o sucesso deste evento e para o desenvolvimento científico dos estudantes.')
    doc.add_paragraph('')
    doc.add_paragraph('Atenciosamente,')
    doc.add_paragraph('Comissão Organizadora da FECITEL')
    
    return doc

def create_premiacao_pdf():
    """Cria o PDF da premiação"""
    temp_dir = create_temp_dir()
    filename = os.path.join(temp_dir, "premiacao.pdf")
    
    c = canvas.Canvas(filename, pagesize=A4)
    width, height = A4
    
    # Título
    c.setFont("Helvetica-Bold", 24)
    c.drawString(width/2 - 80, height - 100, "PREMIAÇÃO FECITEL")
    
    # Categorias
    c.setFont("Helvetica-Bold", 16)
    y_position = height - 200
    c.drawString(100, y_position, "Categorias de Premiação:")
    
    c.setFont("Helvetica", 12)
    y_position -= 30
    c.drawString(120, y_position, "🥇 1º Lugar - Medalha de Ouro + Certificado")
    y_position -= 20
    c.drawString(120, y_position, "🥈 2º Lugar - Medalha de Prata + Certificado")
    y_position -= 20
    c.drawString(120, y_position, "🥉 3º Lugar - Medalha de Bronze + Certificado")
    y_position -= 20
    c.drawString(120, y_position, "🏆 Menção Honrosa - Certificado")
    
    c.save()
    return filename

def create_relacao_trabalhos_pptx():
    """Cria a apresentação com relação de trabalhos"""
    prs = Presentation()
    
    # Slide de título
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "RELATION DE TRABALHOS"
    subtitle.text = "FECITEL - Feira de Ciência e Tecnologia"
    
    # Slide com lista de trabalhos
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Trabalhos Inscritos"
    
    tf = body_shape.text_frame
    tf.text = "Lista de Projetos:"
    
    p = tf.add_paragraph()
    p.text = "Projeto 1: Sistema de Irrigação Automática"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Projeto 2: Energia Solar Sustentável"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Projeto 3: Reciclagem Inteligente"
    p.level = 1
    
    temp_dir = create_temp_dir()
    filename = os.path.join(temp_dir, "relacao_trabalhos.pptx")
    prs.save(filename)
    return filename

def create_script_encerramento_doc():
    """Cria o script de encerramento"""
    doc = Document()
    doc.add_heading('SCRIPT ENCERRAMENTO SCT E FECITEL', 0)
    
    doc.add_heading('Programação do Encerramento', level=1)
    doc.add_paragraph('1. Abertura - Diretor da Instituição')
    doc.add_paragraph('2. Apresentação dos Resultados - Coordenador da FECITEL')
    doc.add_paragraph('3. Premiação dos Vencedores')
    doc.add_paragraph('4. Agradecimentos aos Participantes')
    doc.add_paragraph('5. Encerramento')
    
    doc.add_heading('Fala do Diretor', level=1)
    doc.add_paragraph('"É com grande satisfação que encerramos mais uma edição da FECITEL. Parabenizamos todos os participantes pelo excelente trabalho desenvolvido."')
    
    return doc

def create_slide_fecitel_odp():
    """Cria o slide da FECITEL em formato ODP"""
    doc = OpenDocumentText()
    
    # Adicionar título
    h = text.H(outlinelevel=1, stylename="Heading 1")
    h.addNewText("FECITEL")
    doc.text.addElement(h)
    
    # Adicionar parágrafo
    p = text.P()
    p.addNewText("Feira de Ciência e Tecnologia")
    doc.text.addElement(p)
    
    temp_dir = create_temp_dir()
    filename = os.path.join(temp_dir, "slide_fecitel.odp")
    doc.save(filename)
    return filename

def create_certificado_feiras_doc():
    """Cria o certificado de feiras"""
    doc = Document()
    doc.add_heading('CERTIFICADO', 0)
    doc.add_paragraph('')
    doc.add_paragraph('Certificamos que')
    doc.add_paragraph('')
    doc.add_paragraph('_________________________________')
    doc.add_paragraph('')
    doc.add_paragraph('participou da FECITEL - Feira de Ciência e Tecnologia')
    doc.add_paragraph('')
    doc.add_paragraph(f'Data: {datetime.now().strftime("%d/%m/%Y")}')
    doc.add_paragraph('')
    doc.add_paragraph('Assinatura: _________________')
    
    return doc

def create_ficha_avaliacao_pdf():
    """Cria a ficha de avaliação em PDF"""
    temp_dir = create_temp_dir()
    filename = os.path.join(temp_dir, "ficha_avaliacao.pdf")
    
    c = canvas.Canvas(filename, pagesize=A4)
    width, height = A4
    
    # Título
    c.setFont("Helvetica-Bold", 18)
    c.drawString(width/2 - 80, height - 50, "FICHA DE AVALIAÇÃO FECITEL")
    
    # Informações do projeto
    c.setFont("Helvetica-Bold", 12)
    y_position = height - 100
    c.drawString(50, y_position, "Nome do Projeto: _________________________________")
    y_position -= 20
    c.drawString(50, y_position, "Escola: _________________________________________")
    y_position -= 20
    c.drawString(50, y_position, "Avaliador: _______________________________________")
    
    # Critérios
    y_position -= 40
    c.setFont("Helvetica-Bold", 12)
    c.drawString(50, y_position, "CRITÉRIOS DE AVALIAÇÃO:")
    
    criteria = [
        "Criatividade e Inovação (0-25): _____",
        "Metodologia Científica (0-25): _____",
        "Apresentação e Comunicação (0-20): _____",
        "Relevância Social (0-15): _____",
        "Viabilidade Técnica (0-15): _____"
    ]
    
    y_position -= 30
    c.setFont("Helvetica", 10)
    for criterion in criteria:
        c.drawString(60, y_position, criterion)
        y_position -= 15
    
    # Total
    y_position -= 20
    c.setFont("Helvetica-Bold", 12)
    c.drawString(50, y_position, "TOTAL (0-100): _____")
    
    c.save()
    return filename

@router.get("/anais")
async def generate_anais():
    """Gera o documento Anais FECITEL.doc"""
    try:
        doc = create_anais_doc()
        temp_dir = create_temp_dir()
        filename = os.path.join(temp_dir, "anais_fecitel.docx")
        doc.save(filename)
        return FileResponse(filename, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename="anais_fecitel.docx")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erro ao gerar anais: {str(e)}")

@router.get("/apresentacao")
async def generate_apresentacao():
    """Gera o PDF da apresentação da feira"""
    try:
        filename = create_apresentacao_pdf()
        return FileResponse(filename, media_type="application/pdf", filename="apresentacao_feira.pdf")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erro ao gerar apresentação: {str(e)}")

@router.get("/instrucoes")
async def generate_instrucoes():
    """Gera o documento de instruções para avaliação"""
    try:
        doc = create_instrucoes_doc()
        temp_dir = create_temp_dir()
        filename = os.path.join(temp_dir, "instrucoes_avaliacao.docx")
        doc.save(filename)
        return FileResponse(filename, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename="instrucoes_avaliacao.docx")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erro ao gerar instruções: {str(e)}")

@router.get("/mensagem-avaliador")
async def generate_mensagem_avaliador():
    """Gera a mensagem para o avaliador"""
    try:
        doc = create_mensagem_avaliador_doc()
        temp_dir = create_temp_dir()
        filename = os.path.join(temp_dir, "mensagem_avaliador.docx")
        doc.save(filename)
        return FileResponse(filename, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename="mensagem_avaliador.docx")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erro ao gerar mensagem: {str(e)}")

@router.get("/premiacao")
async def generate_premiacao():
    """Gera o PDF da premiação"""
    try:
        filename = create_premiacao_pdf()
        return FileResponse(filename, media_type="application/pdf", filename="premiacao.pdf")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erro ao gerar premiação: {str(e)}")

@router.get("/relacao-trabalhos")
async def generate_relacao_trabalhos():
    """Gera a apresentação com relação de trabalhos"""
    try:
        filename = create_relacao_trabalhos_pptx()
        return FileResponse(filename, media_type="application/vnd.openxmlformats-presentationml.presentation", filename="relacao_trabalhos.pptx")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erro ao gerar relação de trabalhos: {str(e)}")

@router.get("/script-encerramento")
async def generate_script_encerramento():
    """Gera o script de encerramento"""
    try:
        doc = create_script_encerramento_doc()
        temp_dir = create_temp_dir()
        filename = os.path.join(temp_dir, "script_encerramento.docx")
        doc.save(filename)
        return FileResponse(filename, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename="script_encerramento.docx")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erro ao gerar script: {str(e)}")

@router.get("/slide-fecitel")
async def generate_slide_fecitel():
    """Gera o slide da FECITEL"""
    try:
        filename = create_slide_fecitel_odp()
        return FileResponse(filename, media_type="application/vnd.oasis.opendocument.text", filename="slide_fecitel.odp")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erro ao gerar slide: {str(e)}")

@router.get("/certificado")
async def generate_certificado():
    """Gera o certificado de feiras"""
    try:
        doc = create_certificado_feiras_doc()
        temp_dir = create_temp_dir()
        filename = os.path.join(temp_dir, "certificado_feiras.docx")
        doc.save(filename)
        return FileResponse(filename, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename="certificado_feiras.docx")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erro ao gerar certificado: {str(e)}")

@router.get("/ficha-avaliacao")
async def generate_ficha_avaliacao():
    """Gera a ficha de avaliação em PDF"""
    try:
        filename = create_ficha_avaliacao_pdf()
        return FileResponse(filename, media_type="application/pdf", filename="ficha_avaliacao.pdf")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erro ao gerar ficha de avaliação: {str(e)}") 