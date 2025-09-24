from fastapi import APIRouter, Depends, HTTPException, status, Query
from fastapi.responses import FileResponse
from sqlalchemy.orm import Session
from app.database import get_db
from app.models.document import Document
from app.schemas.documents import (
    DocumentCreate, DocumentUpdate, DocumentListResponse, DocumentDetailResponse, DocumentGenerationRequest
)
import os
import tempfile
from datetime import datetime
from docx import Document as DocxDocument
from docx.shared import Inches
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.units import inch
from pptx import Presentation
from pptx.util import Inches
from odf import text, teletype
from odf.opendocument import OpenDocumentText

router = APIRouter()

@router.get("/", response_model=DocumentListResponse)
async def get_documents(
    skip: int = Query(0, ge=0),
    limit: int = Query(100, ge=1, le=1000),
    db: Session = Depends(get_db)
):
    try:
        documents = db.query(Document).filter(Document.deleted_at == None).offset(skip).limit(limit).all()
        
        document_data = []
        for document in documents:
            document_dict = {
                "id": document.id,
                "name": document.name,
                "type": document.type,
                "header": document.header,
                "logo1": document.logo1,
                "logo2": document.logo2,
                "content": document.content,
                "created_at": document.created_at,
                "updated_at": document.updated_at,
                "deleted_at": document.deleted_at
            }
            document_data.append(document_dict)
        
        return DocumentListResponse(
            status=True,
            message="Documents retrieved successfully",
            data=document_data
        )
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Error retrieving documents: {str(e)}"
        )

@router.get("/{document_id}", response_model=DocumentDetailResponse)
async def get_document(
    document_id: int,
    db: Session = Depends(get_db)
):
    try:
        document = db.query(Document).filter(Document.id == document_id, Document.deleted_at == None).first()
        
        if not document:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="Document not found"
            )
        
        document_dict = {
            "id": document.id,
            "name": document.name,
            "type": document.type,
            "header": document.header,
            "logo1": document.logo1,
            "logo2": document.logo2,
            "content": document.content,
            "created_at": document.created_at,
            "updated_at": document.updated_at,
            "deleted_at": document.deleted_at
        }
        
        return DocumentDetailResponse(
            status=True,
            message="Document retrieved successfully",
            data=document_dict
        )
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Error retrieving document: {str(e)}"
        )

@router.post("/", response_model=DocumentDetailResponse)
async def create_document(document_data: DocumentCreate, db: Session = Depends(get_db)):
    try:
        document = Document(
            name=document_data.name,
            type=document_data.type,
            header=document_data.header,
            logo1=document_data.logo1,
            logo2=document_data.logo2,
            content=document_data.content
        )
        
        db.add(document)
        db.commit()
        db.refresh(document)
        
        document_dict = {
            "id": document.id,
            "name": document.name,
            "type": document.type,
            "header": document.header,
            "logo1": document.logo1,
            "logo2": document.logo2,
            "content": document.content,
            "created_at": document.created_at,
            "updated_at": document.updated_at,
            "deleted_at": document.deleted_at
        }
        
        return DocumentDetailResponse(
            status=True,
            message="Document created successfully",
            data=document_dict
        )
    except Exception as e:
        db.rollback()
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Error creating document: {str(e)}"
        )

@router.put("/{document_id}", response_model=DocumentDetailResponse)
async def update_document(
    document_id: int,
    document_data: DocumentUpdate,
    db: Session = Depends(get_db)
):
    try:
        document = db.query(Document).filter(Document.id == document_id, Document.deleted_at == None).first()
        
        if not document:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="Document not found"
            )
        
        update_data = document_data.dict(exclude_unset=True)
        for field, value in update_data.items():
            setattr(document, field, value)
        
        db.commit()
        db.refresh(document)
        
        document_dict = {
            "id": document.id,
            "name": document.name,
            "type": document.type,
            "header": document.header,
            "logo1": document.logo1,
            "logo2": document.logo2,
            "content": document.content,
            "created_at": document.created_at,
            "updated_at": document.updated_at,
            "deleted_at": document.deleted_at
        }
        
        return DocumentDetailResponse(
            status=True,
            message="Document updated successfully",
            data=document_dict
        )
    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Error updating document: {str(e)}"
        )

@router.delete("/{document_id}")
async def delete_document(document_id: int, db: Session = Depends(get_db)):
    try:
        document = db.query(Document).filter(Document.id == document_id, Document.deleted_at == None).first()
        
        if not document:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail="Document not found"
            )
        
        from datetime import datetime
        document.deleted_at = datetime.utcnow()
        
        db.commit()
        
        return {
            "status": True,
            "message": "Document deleted successfully"
        }
    except HTTPException:
        raise
    except Exception as e:
        db.rollback()
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Error deleting document: {str(e)}"
        )

def create_temp_dir():
    temp_dir = tempfile.mkdtemp(prefix="fecitel_docs_")
    return temp_dir

def create_anais_doc():
    doc = DocxDocument()
    doc.add_heading('ANAIS FECITEL', 0)
    doc.add_paragraph('Feira de Ciência e Tecnologia')
    doc.add_paragraph(f'Data: {datetime.now().strftime("%d/%m/%Y")}')
    
    doc.add_heading('Sumário', level=1)
    doc.add_paragraph('1. Apresentação')
    doc.add_paragraph('2. Trabalhos Apresentados')
    doc.add_paragraph('3. Avaliações')
    doc.add_paragraph('4. Premiação')
    
    doc.add_heading('1. Apresentação', level=1)
    doc.add_paragraph('Este documento apresenta os anais da Feira de Ciência e Tecnologia (FECITEL), contendo todos os trabalhos apresentados, suas avaliações e resultados finais.')
    
    return doc

def create_participacao_pdf():
    temp_dir = create_temp_dir()
    filename = os.path.join(temp_dir, "participacao_feira.pdf")
    
    c = canvas.Canvas(filename, pagesize=A4)
    width, height = A4
    
    c.setFont("Helvetica-Bold", 24)
    c.drawString(width/2 - 100, height - 100, "FECITEL")
    c.setFont("Helvetica", 16)
    c.drawString(width/2 - 80, height - 130, "Feira de Ciência e Tecnologia")
    
    c.setFont("Helvetica", 12)
    y_position = height - 200
    c.drawString(100, y_position, "• Evento anual de ciência e tecnologia")
    y_position -= 20
    c.drawString(100, y_position, "• Participação de escolas da região")
    y_position -= 20
    c.drawString(100, y_position, "• Avaliação por comissão especializada")
    y_position -= 20
    c.drawString(100, y_position, f"• Data: {datetime.now().strftime('%d/%m/%Y')}")
    
    c.save()
    return filename

def create_instrucoes_doc():
    doc = DocxDocument()
    doc.add_heading('INSTRUÇÕES PARA AVALIAÇÃO DE TRABALHOS NA FECITEL 2024', 0)
    
    # Introdução
    doc.add_paragraph('Prezado avaliador, o presente documento traz orientações importantes para o')
    doc.add_paragraph('sucesso da sua avaliação. Peço que o leia com atenção e em caso de dúvidas, peça')
    doc.add_paragraph('ajuda a um membro da comissão organizadora.')
    doc.add_paragraph('')
    
    # Item 1
    doc.add_heading('1. ITENS DE AVALIAÇÃO OBRIGATÓRIA', level=1)
    doc.add_paragraph('Os estudantes deverão realizar uma exposição oral do trabalho na qual serão')
    doc.add_paragraph('analisados os seguintes itens:')
    doc.add_paragraph('a) Resumo expandido: está disponível na pasta de rede "Resumos Fecitel" (O acesso a')
    doc.add_paragraph('pasta será explicado no item 2);')
    doc.add_paragraph('b) Banner: disponível na sala de apresentação;')
    doc.add_paragraph('c) Caderno de Campo ou Diário de Bordo: deve ser apresentado pelos estudantes e')
    doc.add_paragraph('d) Relatório do trabalho: deve ser apresentado pelos estudantes.')
    doc.add_paragraph('')
    doc.add_paragraph('A Comissão Organizadora já realizou uma pré-avaliação dos trabalhos, portanto, seu')
    doc.add_paragraph('julgamento tem relação com o mérito científico e relevância do trabalho conforme itens de')
    doc.add_paragraph('avaliação disponíveis no Sistema de Avaliação (item 4).')
    doc.add_paragraph('')
    doc.add_paragraph('Caso o estudante não apresente o Caderno de Campo ou Diário de Bordo, bem como')
    doc.add_paragraph('o Relatório de Trabalho na apresentação, a pontuação referente a esses itens deve ser')
    doc.add_paragraph('zerada no Sistema de Avaliação.')
    doc.add_paragraph('')
    doc.add_paragraph('Cada trabalho é avaliado por três avaliadores de forma independente, sendo assim, a')
    doc.add_paragraph('nota atribuída por você será somada a outras duas para o cálculo da pontuação final.')
    doc.add_paragraph('')
    
    # Item 2
    doc.add_heading('2. COMO ACESSAR OS COMPUTADORES E A PASTA DE RESUMOS?', level=1)
    doc.add_paragraph('a) Ligar o computador.')
    doc.add_paragraph('b) Clicar em "Opções de Entrada" em seguida clicar no símbolo de chave.')
    doc.add_paragraph('c) Logar utilizando as seguintes credenciais: Usuário: aluno Senha: alunoifms')
    doc.add_paragraph('d) Abrir o "Explorador de Arquivos" e digitar na barra de endereços: \\\\10.8.32.3\\arquivos\\')
    doc.add_paragraph('e) Acessar a pasta "Resumos Fecitel".')
    doc.add_paragraph('f) Acessar a pasta da Área do trabalho.')
    doc.add_paragraph('g) Localizar o(s) trabalho(s) designado(s) para avaliação.')
    doc.add_paragraph('')
    
    # Item 3
    doc.add_heading('3. COMO SABER QUAIS TRABALHOS TENHO QUE AVALIAR?', level=1)
    doc.add_paragraph('Cada avaliador tem de 2 a 5 trabalhos para ser avaliado, dependendo da área de')
    doc.add_paragraph('formação e do currículo profissional.')
    doc.add_paragraph('Acessando o Sistema de Avaliação com seu PIN (item 4) os trabalhos disponíveis para')
    doc.add_paragraph('sua avaliação aparecerão na tela.')
    doc.add_paragraph('Você também pode consultar os trabalhos a serem avaliados na Planilha colada na')
    doc.add_paragraph('Sala de Avaliações.')
    doc.add_paragraph('')
    
    # Item 4
    doc.add_heading('4. COMO ACESSAR O SISTEMA DE AVALIAÇÃO E AVALIAR OS TRABALHOS?', level=1)
    doc.add_paragraph('Você poderá acessar o sistema de avaliação utilizando o computador ou seu')
    doc.add_paragraph('smartphone, conforme for mais conveniente.')
    doc.add_paragraph('Todos os critérios de avaliação que foram enviados para conhecimento via e-mail,')
    doc.add_paragraph('estão cadastrados no Sistema de Avaliação.')
    doc.add_paragraph('')
    doc.add_paragraph('a) Pelo computador: Abra o navegador de internet e digite o endereço http://10.8.16.52')
    doc.add_paragraph('b) Pelo smartphone: Aponte a câmera do seu smartphone para o QRCode:')
    doc.add_paragraph('')
    doc.add_paragraph('Digite seu PIN.')
    doc.add_paragraph('Os trabalhos disponíveis para sua avaliação estarão visíveis')
    doc.add_paragraph('na tela principal.')
    doc.add_paragraph('Ao clicar no trabalho, cada critério de avaliação será aberto')
    doc.add_paragraph('para que você clique na nota desejada. Após atribuir a nota,')
    doc.add_paragraph('clique em avançar.')
    doc.add_paragraph('Ao final da avaliação, você poderá ver um resumo da')
    doc.add_paragraph('avaliação. Clique em "Finalizar Avaliação"')
    doc.add_paragraph('')
    doc.add_paragraph('Utilize nossa rede Wi-Fi IFMS-ADMINISTRATIVO')
    doc.add_paragraph('Usuário: visitantes')
    doc.add_paragraph('Senha: Ifms.2024')
    doc.add_paragraph('')
    doc.add_paragraph('Quando finalizar sua avaliação, não se esqueça de assinar a lista de presença e')
    doc.add_paragraph('pegar sua lembrança.')
    doc.add_paragraph('Enviaremos o Certificado de Avaliador para o e-mail cadastrado.')
    doc.add_paragraph('Qualquer dúvida, chame um membro da Comissão Organizadora para te auxiliar.')
    doc.add_paragraph('Agradecemos sua participação!')
    doc.add_paragraph('')
    doc.add_paragraph('Atenciosamente,')
    doc.add_paragraph('Comissão Organizadora da Fecitel 2024')
    
    return doc

def create_mensagem_avaliador_doc():
    doc = DocxDocument()
    doc.add_heading('MENSAGEM AO AVALIADOR', 0)
    
    doc.add_paragraph('Prezado(a) Avaliador(a),')
    doc.add_paragraph('')
    doc.add_paragraph('Agradecemos sua participação na FECITEL - Feira de Ciência e Tecnologia.')
    doc.add_paragraph('')
    doc.add_paragraph('Sua avaliação é fundamental para o sucesso deste evento e para o desenvolvimento científico dos estudantes.')
    doc.add_paragraph('')
    doc.add_paragraph('Atenciosamente,')
    doc.add_paragraph('Comissão Organizadora da FECITEL')
    
    return doc

def create_premiacao_pdf():
    temp_dir = create_temp_dir()
    filename = os.path.join(temp_dir, "premiacao.pdf")
    
    c = canvas.Canvas(filename, pagesize=A4)
    width, height = A4
    
    c.setFont("Helvetica-Bold", 24)
    c.drawString(width/2 - 80, height - 100, "PREMIAÇÃO FECITEL")
    
    c.setFont("Helvetica-Bold", 16)
    y_position = height - 200
    c.drawString(100, y_position, "Categorias de Premiação:")
    
    c.setFont("Helvetica", 12)
    y_position -= 30
    c.drawString(120, y_position, "🥇 1º Lugar - Medalha de Ouro + Certificado")
    y_position -= 20
    c.drawString(120, y_position, "🥈 2º Lugar - Medalha de Prata + Certificado")
    y_position -= 20
    c.drawString(120, y_position, "🥉 3º Lugar - Medalha de Bronze + Certificado")
    y_position -= 20
    c.drawString(120, y_position, "🏆 Menção Honrosa - Certificado")
    
    c.save()
    return filename

def create_relacao_trabalhos_pptx():
    prs = Presentation()
    
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "RELATION DE TRABALHOS"
    subtitle.text = "FECITEL - Feira de Ciência e Tecnologia"
    
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Trabalhos Inscritos"
    
    tf = body_shape.text_frame
    tf.text = "Lista de Projetos:"
    
    p = tf.add_paragraph()
    p.text = "Projeto 1: Sistema de Irrigação Automática"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Projeto 2: Energia Solar Sustentável"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Projeto 3: Reciclagem Inteligente"
    p.level = 1
    
    temp_dir = create_temp_dir()
    filename = os.path.join(temp_dir, "relacao_trabalhos.pptx")
    prs.save(filename)
    return filename

def create_script_encerramento_doc():
    doc = DocxDocument()
    doc.add_heading('SCRIPT ENCERRAMENTO SCT E FECITEL', 0)
    
    doc.add_heading('Programação do Encerramento', level=1)
    doc.add_paragraph('1. Abertura - Diretor da Instituição')
    doc.add_paragraph('2. Apresentação dos Resultados - Coordenador da FECITEL')
    doc.add_paragraph('3. Premiação dos Vencedores')
    doc.add_paragraph('4. Agradecimentos aos Participantes')
    doc.add_paragraph('5. Encerramento')
    
    doc.add_heading('Fala do Diretor', level=1)
    doc.add_paragraph('"É com grande satisfação que encerramos mais uma edição da FECITEL. Parabenizamos todos os participantes pelo excelente trabalho desenvolvido."')
    
    return doc

def create_slide_fecitel_odp():
    doc = OpenDocumentText()
    
    h = text.H(outlinelevel=1, stylename="Heading 1")
    h.addNewText("FECITEL")
    doc.text.addElement(h)
    
    p = text.P()
    p.addNewText("Feira de Ciência e Tecnologia")
    doc.text.addElement(p)
    
    temp_dir = create_temp_dir()
    filename = os.path.join(temp_dir, "slide_fecitel.odp")
    doc.save(filename)
    return filename

def create_certificado_feiras_doc():
    doc = DocxDocument()
    doc.add_heading('CERTIFICADO', 0)
    doc.add_paragraph('')
    doc.add_paragraph('Certificamos que')
    doc.add_paragraph('')
    doc.add_paragraph('_________________________________')
    doc.add_paragraph('')
    doc.add_paragraph('participou da FECITEL - Feira de Ciência e Tecnologia')
    doc.add_paragraph('')
    doc.add_paragraph(f'Data: {datetime.now().strftime("%d/%m/%Y")}')
    doc.add_paragraph('')
    doc.add_paragraph('Assinatura: _________________')
    
    return doc

def create_ficha_avaliacao_pdf():
    temp_dir = create_temp_dir()
    filename = os.path.join(temp_dir, "ficha_avaliacao.pdf")
    
    c = canvas.Canvas(filename, pagesize=A4)
    width, height = A4
    
    c.setFont("Helvetica-Bold", 18)
    c.drawString(width/2 - 80, height - 50, "FICHA DE AVALIAÇÃO FECITEL")
    
    c.setFont("Helvetica-Bold", 12)
    y_position = height - 100
    c.drawString(50, y_position, "Nome do Projeto: _________________________________")
    y_position -= 20
    c.drawString(50, y_position, "Escola: _________________________________________")
    y_position -= 20
    c.drawString(50, y_position, "Avaliador: _______________________________________")
    
    y_position -= 40
    c.setFont("Helvetica-Bold", 12)
    c.drawString(50, y_position, "CRITÉRIOS DE AVALIAÇÃO:")
    
    criteria = [
        "Criatividade e Inovação (0-25): _____",
        "Metodologia Científica (0-25): _____",
        "Apresentação e Comunicação (0-20): _____",
        "Relevância Social (0-15): _____",
        "Viabilidade Técnica (0-15): _____"
    ]
    
    y_position -= 30
    c.setFont("Helvetica", 10)
    for criterion in criteria:
        c.drawString(60, y_position, criterion)
        y_position -= 15
    
    y_position -= 20
    c.setFont("Helvetica-Bold", 12)
    c.drawString(50, y_position, "TOTAL (0-100): _____")
    
    c.save()
    return filename

# @router.get("/generate/anais")
# async def generate_anais():
#     try:
#         doc = create_anais_doc()
#         temp_dir = create_temp_dir()
#         filename = os.path.join(temp_dir, "anais_fecitel.docx")
#         doc.save(filename)
#         return FileResponse(filename, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename="anais_fecitel.docx")
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"Erro ao gerar anais: {str(e)}")

@router.post("/generate/participacao")
async def generate_participacao(request: DocumentGenerationRequest, db: Session = Depends(get_db)):
    try:        
        filename = create_participacao_pdf()
        return FileResponse(filename, media_type="application/pdf", filename="participacao_feira.pdf")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erro ao gerar apresentação: {str(e)}")

@router.post("/generate/instrucoes")
async def generate_instrucoes(request: DocumentGenerationRequest, db: Session = Depends(get_db)):
    try:
        doc = create_instrucoes_doc()
        temp_dir = create_temp_dir()
        filename = os.path.join(temp_dir, "instrucoes_avaliacao.docx")
        doc.save(filename)
        return FileResponse(filename, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename="instrucoes_avaliacao.docx")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erro ao gerar instruções: {str(e)}")

# @router.get("/generate/mensagem-avaliador")
# async def generate_mensagem_avaliador():
#     try:
#         doc = create_mensagem_avaliador_doc()
#         temp_dir = create_temp_dir()
#         filename = os.path.join(temp_dir, "mensagem_avaliador.docx")
#         doc.save(filename)
#         return FileResponse(filename, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename="mensagem_avaliador.docx")
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"Erro ao gerar mensagem: {str(e)}")

@router.post("/generate/premiacao")
async def generate_premiacao(request: DocumentGenerationRequest, db: Session = Depends(get_db)):
    try:
        filename = create_premiacao_pdf()
        return FileResponse(filename, media_type="application/pdf", filename="premiacao.pdf")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Erro ao gerar premiação: {str(e)}")

# @router.get("/generate/relacao-trabalhos")
# async def generate_relacao_trabalhos():
#     try:
#         filename = create_relacao_trabalhos_pptx()
#         return FileResponse(filename, media_type="application/vnd.openxmlformats-presentationml.presentation", filename="relacao_trabalhos.pptx")
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"Erro ao gerar relação de trabalhos: {str(e)}")

# @router.get("/generate/script-encerramento")
# async def generate_script_encerramento():
#     try:
#         doc = create_script_encerramento_doc()
#         temp_dir = create_temp_dir()
#         filename = os.path.join(temp_dir, "script_encerramento.docx")
#         doc.save(filename)
#         return FileResponse(filename, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename="script_encerramento.docx")
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"Erro ao gerar script: {str(e)}")

# @router.get("/generate/slide-fecitel")
# async def generate_slide_fecitel():
#     try:
#         filename = create_slide_fecitel_odp()
#         return FileResponse(filename, media_type="application/vnd.oasis.opendocument.text", filename="slide_fecitel.odp")
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"Erro ao gerar slide: {str(e)}")

# @router.get("/generate/certificado")
# async def generate_certificado():
#     try:
#         doc = create_certificado_feiras_doc()
#         temp_dir = create_temp_dir()
#         filename = os.path.join(temp_dir, "certificado_feiras.docx")
#         doc.save(filename)
#         return FileResponse(filename, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename="certificado_feiras.docx")
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"Erro ao gerar certificado: {str(e)}")

# @router.get("/generate/ficha-avaliacao")
# async def generate_ficha_avaliacao():
#     try:
#         filename = create_ficha_avaliacao_pdf()
#         return FileResponse(filename, media_type="application/pdf", filename="ficha_avaliacao.pdf")
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=f"Erro ao gerar ficha de avaliação: {str(e)}") 